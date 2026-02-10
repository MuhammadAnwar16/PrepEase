import os
from typing import Literal

import fitz
from pptx import Presentation

FileType = Literal["pdf", "pptx"]


def _detect_file_type(file_path: str) -> FileType:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return "pdf"
    if ext in {".pptx", ".ppt"}:
        return "pptx"
    raise ValueError("Unsupported file type. Only PDF or PPTX is allowed.")


def extract_text(file_path: str) -> str:
    file_type = _detect_file_type(file_path)

    if file_type == "pdf":
        text_parts = []
        with fitz.open(file_path) as doc:
            for page in doc:
                text_parts.append(page.get_text())
        return "\n".join(text_parts).strip()

    presentation = Presentation(file_path)
    text_runs = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text_runs.append(shape.text)
    return "\n".join(text_runs).strip()
